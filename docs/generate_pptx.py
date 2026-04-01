#!/usr/bin/env python3
"""
Generate POLLY Platform Overview PowerPoint presentation.

Uses screenshots from static/guide/ for slide visuals.
Output: docs/POLLY_Platform_Overview.pptx

Usage:
    python docs/generate_pptx.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).parent.parent
GUIDE_DIR = ROOT / "static" / "guide"
OUTPUT = Path(__file__).parent / "POLLY_Platform_Overview.pptx"

# Brand colors
POLLY_PRIMARY = RGBColor(0x63, 0x66, 0xF1)   # Indigo
POLLY_SECONDARY = RGBColor(0x06, 0xB6, 0xD4)  # Cyan
POLLY_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Dark navy
POLLY_CARD = RGBColor(0x1E, 0x29, 0x3B)       # Card bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, font_size=18,
              bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_image_safe(slide, img_name, left, top, width):
    img_path = GUIDE_DIR / img_name
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), Inches(width))
        return True
    return False


def _title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_slide_bg(slide, POLLY_DARK)
    _add_text(slide, 1, 1.5, 8, 1, "POLLY", 44, bold=True, color=POLLY_PRIMARY, alignment=PP_ALIGN.CENTER)
    _add_text(slide, 1, 2.5, 8, 0.8, "AI Marketing Platform for Financial Advisors", 24, color=POLLY_SECONDARY, alignment=PP_ALIGN.CENTER)
    _add_text(slide, 1.5, 4, 7, 1.5,
              "Compliant campaign creation, multi-channel distribution, "
              "and real-time analytics — powered by AI agents.",
              16, color=MUTED, alignment=PP_ALIGN.CENTER)
    _add_text(slide, 3, 6.5, 4, 0.5, "predictivelabs.ai", 12, color=MUTED, alignment=PP_ALIGN.CENTER)


def _section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, POLLY_CARD)
    _add_text(slide, 1, 2.5, 8, 1, title, 36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        _add_text(slide, 1, 3.8, 8, 1, subtitle, 16, color=MUTED, alignment=PP_ALIGN.CENTER)


def _content_slide(prs, title, bullets, img_name=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, POLLY_DARK)
    _add_text(slide, 0.5, 0.3, 9, 0.6, title, 24, bold=True, color=POLLY_PRIMARY)

    if img_name and (GUIDE_DIR / img_name).exists():
        # Left: bullets, Right: image
        bullet_width = 4.5
        y = 1.2
        for bullet in bullets:
            _add_text(slide, 0.7, y, bullet_width, 0.4, f"  {bullet}", 13, color=WHITE)
            y += 0.45
        _add_image_safe(slide, img_name, 5.3, 1.0, 4.5)
    else:
        y = 1.3
        for bullet in bullets:
            _add_text(slide, 0.7, y, 8.5, 0.45, f"  {bullet}", 14, color=WHITE)
            y += 0.5


def _two_image_slide(prs, title, img1, cap1, img2, cap2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, POLLY_DARK)
    _add_text(slide, 0.5, 0.3, 9, 0.6, title, 24, bold=True, color=POLLY_PRIMARY)
    _add_image_safe(slide, img1, 0.3, 1.1, 4.6)
    _add_text(slide, 0.3, 5.8, 4.6, 0.4, cap1, 10, color=MUTED, alignment=PP_ALIGN.CENTER)
    _add_image_safe(slide, img2, 5.1, 1.1, 4.6)
    _add_text(slide, 5.1, 5.8, 4.6, 0.4, cap2, 10, color=MUTED, alignment=PP_ALIGN.CENTER)


def generate():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    _title_slide(prs)

    # Slide 2: Problem Statement
    _content_slide(prs, "The Challenge", [
        "Financial advisors spend 2-3 FTEs on campaign management",
        "Compliance review is manual and error-prone",
        "Multi-channel coordination is fragmented",
        "No unified analytics across WhatsApp, email, social, CRM",
        "Marketing content must comply with MiFID II, PRIIPs, FCA",
    ])

    # Slide 3: Solution
    _content_slide(prs, "POLLY — Your AI Marketing Team", [
        "9 AI agents with 54 specialized marketing tools",
        "Built-in MiFID II / PRIIPs / FCA compliance guardrails",
        "Multi-channel: WhatsApp, Telegram, email, LinkedIn, X, Instagram, TikTok",
        "Campaign automation: A/B testing, follow-ups, lead scoring",
        "Conversational interface — just tell POLLY what you need",
    ], "01_home.png")

    # Slide 4: Platform Overview
    _section_slide(prs, "Platform Overview", "9 agents, 54 tools, 7+ channels")

    # Slide 5: Chat Interface
    _two_image_slide(prs, "Conversational AI Interface",
                     "07_chat_starter.png", "6 starter skill buttons",
                     "08_chat_campaign.png", "Campaign creation via chat")

    # Slide 6: Compliance
    _content_slide(prs, "Compliance Built-In", [
        "AI-powered review against MiFID II, PRIIPs, FCA regulations",
        "Automated risk warnings per jurisdiction (UK, EU, US, APAC)",
        "Target market assessment and negative target definition",
        "Document approval workflow with four-eye management checks",
        "Compliance-approved document set drives all campaigns",
    ], "12_instructions_top.png")

    # Slide 7: Campaign Management
    _content_slide(prs, "Campaign Lifecycle Management", [
        "Warmup campaigns to test market appetite",
        "AI-generated compliant copy for every channel",
        "A/B testing with statistical framework",
        "Automated follow-ups for non-responders",
        "Lead scoring: hot, warm, questions, removal",
        "CRM integration for pipeline tracking",
    ], "09_chat_analytics.png")

    # Slide 8: Multi-Channel
    _content_slide(prs, "Multi-Channel Distribution", [
        "WhatsApp Business — deliver, read receipts, replies",
        "Telegram Bot — messaging and group monitoring",
        "Email — opens, clicks, replies, unsubscribes",
        "LinkedIn & X — social engagement and posting",
        "Instagram & TikTok — reach and engagement analytics",
        "CRM — pipeline movement and lead tracking",
    ], "02_about_top.png")

    # Slide 9: Interactive Demo
    _two_image_slide(prs, "Interactive Demo — Device Simulator",
                     "04_demo.png", "WhatsApp simulator",
                     "15_demo_telegram.png", "Telegram simulator")

    # Slide 10: Demo Analytics
    _two_image_slide(prs, "Demo — Analytics & Campaign Preview",
                     "16_demo_analytics.png", "Live analytics dashboard",
                     "17_demo_campaign_preview.png", "A/B variant preview")

    # Slide 11: Profile & Integrations
    _two_image_slide(prs, "Profile & API Integrations",
                     "10_profile_top.png", "Integration status dashboard",
                     "11_profile_skills.png", "39 marketing skills")

    # Slide 12: Instructions Editor
    _content_slide(prs, "Customizable Agent Instructions", [
        "Edit system prompts for each of the 9 AI agents",
        "Per-user customization with global defaults",
        "Admin-only global instructions prepended to all agents",
        "Template variables: {{today}} for dynamic date insertion",
        "Reload from database without restart",
        "Prompts stored in PostgreSQL (polly.prompts table)",
    ], "12_instructions_top.png")

    # Slide 13: Architecture
    _content_slide(prs, "Technical Architecture", [
        "Python 3.13 + FastHTML web framework",
        "PostgreSQL database (10 tables in polly schema)",
        "XAI/Grok LLM for content generation",
        "Arcade.dev for social media posting",
        "Playwright for page analysis (CRO/SEO)",
        "Docker deployment via Coolify",
        "Session-based auth with bcrypt password hashing",
    ])

    # Slide 14: Agent Overview
    _content_slide(prs, "9 AI Agents — 54 Tools", [
        "Content (7) — FAQs, teasers, pitch decks, copywriting",
        "Strategy (6) — market research, competitor analysis, backtesting",
        "Compliance (6) — regulatory review, document approval, risk warnings",
        "Campaign (7) — creation, workflows, A/B testing, lead management",
        "Channels (9) — multi-channel monitoring and analytics",
        "Social (3) — post to X, LinkedIn, WhatsApp, Telegram",
        "CRO (8) — conversion optimization, signup flow, forms",
        "SEO (5) — technical audits, schema markup, AI SEO",
        "Ads (3) — paid advertising, analytics tracking",
    ])

    # Slide 15: Deployment
    _content_slide(prs, "Deployment & Security", [
        "Docker container deployed via Coolify",
        "Domain: polly.predictivelabs.ai",
        "PostgreSQL on dedicated server",
        "API keys encrypted, never stored in code",
        "GDPR: removal requests flagged, CRM responsibility",
        "Comprehensive test suite: 52 tests, 135 assertions",
    ])

    # Slide 16: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, POLLY_DARK)
    _add_text(slide, 1, 2, 8, 1, "Ready to meet POLLY?", 36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_text(slide, 1, 3.5, 8, 1, "polly.predictivelabs.ai", 24, color=POLLY_SECONDARY, alignment=PP_ALIGN.CENTER)
    _add_text(slide, 1, 5, 8, 1, "Predictive Labs AI", 16, color=MUTED, alignment=PP_ALIGN.CENTER)

    prs.save(str(OUTPUT))
    print(f"Generated: {OUTPUT} ({OUTPUT.stat().st_size // 1024}K, {len(prs.slides)} slides)")


if __name__ == "__main__":
    generate()
