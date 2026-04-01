#!/usr/bin/env python3
"""
POLLY RAG Pipeline — Process documents into vector embeddings.

Scans doc-data/ for PDF/PPTX/DOCX files, extracts text to markdown,
chunks it, generates embeddings with fastembed, and stores in pgvector.

Usage:
    python tasks/create_rag.py
    python tasks/create_rag.py --doc-dir /path/to/docs
"""
import json
import os
import sys
import time
from pathlib import Path

from dotenv import load_dotenv

ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT))
load_dotenv(ROOT / ".env")


# ---------------------------------------------------------------------------
# Document text extraction
# ---------------------------------------------------------------------------

class DocumentProcessor:
    """Extract markdown text from PDF, PPTX, and DOCX files."""

    def extract(self, filepath: Path) -> str:
        ext = filepath.suffix.lower()
        if ext == ".pdf":
            return self.extract_pdf(filepath)
        elif ext == ".pptx":
            return self.extract_pptx(filepath)
        elif ext == ".docx":
            return self.extract_docx(filepath)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def extract_pdf(self, filepath: Path) -> str:
        try:
            import pymupdf4llm
            return pymupdf4llm.to_markdown(str(filepath))
        except Exception:
            import fitz
            doc = fitz.open(str(filepath))
            pages = []
            for i, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    pages.append(f"## Page {i + 1}\n\n{text}")
            return "\n\n".join(pages)

    def extract_pptx(self, filepath: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)
                    if shape.shape_type == 13 and not title:  # Title placeholder
                        title = shape.text_frame.text.strip()
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip("| "):
                            texts.append(row_text)
            if texts:
                header = f"## Slide {i + 1}" + (f": {title}" if title else "")
                slides.append(f"{header}\n\n" + "\n".join(texts))
        return "\n\n".join(slides)

    def extract_docx(self, filepath: Path) -> str:
        from docx import Document
        doc = Document(str(filepath))
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name.lower() if para.style else ""
            if "heading 1" in style:
                parts.append(f"# {text}")
            elif "heading 2" in style:
                parts.append(f"## {text}")
            elif "heading 3" in style:
                parts.append(f"### {text}")
            else:
                parts.append(text)
        # Extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                rows.append(row_text)
            if rows:
                parts.append("\n".join(rows))
        return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Text chunking
# ---------------------------------------------------------------------------

class TextChunker:
    """Split text into overlapping chunks respecting markdown boundaries."""

    def __init__(self, chunk_size: int = 800, chunk_overlap: int = 150):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def chunk(self, text: str, metadata: dict = None) -> list[dict]:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            separators=["\n## ", "\n### ", "\n\n", "\n", ". ", " "],
            length_function=len,
        )
        docs = splitter.create_documents([text])
        chunks = []
        for i, doc in enumerate(docs):
            chunks.append({
                "chunk_index": i,
                "content": doc.page_content,
                "token_count": len(doc.page_content.split()),
                "metadata": metadata or {},
            })
        return chunks


# ---------------------------------------------------------------------------
# Embedding service
# ---------------------------------------------------------------------------

class EmbeddingService:
    """Generate vector embeddings using fastembed."""

    MODEL = "BAAI/bge-small-en-v1.5"
    DIMENSIONS = 384

    def __init__(self):
        from fastembed import TextEmbedding
        print(f"  Loading embedding model: {self.MODEL}...")
        self._model = TextEmbedding(model_name=self.MODEL)

    def embed(self, texts: list[str]) -> list[list[float]]:
        embeddings = list(self._model.embed(texts))
        return [e.tolist() for e in embeddings]

    @property
    def dimensions(self) -> int:
        return self.DIMENSIONS


# ---------------------------------------------------------------------------
# RAG Indexer — orchestrates the full pipeline
# ---------------------------------------------------------------------------

class RAGIndexer:
    """Process documents: extract → chunk → embed → store in pgvector."""

    def __init__(self, doc_dir: Path):
        from utils.db_pool import DatabasePool
        self.doc_dir = doc_dir
        self.pool = DatabasePool.get()
        self.processor = DocumentProcessor()
        self.chunker = TextChunker()
        self.embedder = EmbeddingService()

    def process_all(self) -> dict:
        files = sorted(
            f for f in self.doc_dir.iterdir()
            if f.suffix.lower() in (".pdf", ".pptx", ".docx")
        )
        print(f"\nFound {len(files)} documents in {self.doc_dir}\n")

        stats = {"documents": [], "total_chunks": 0, "total_files": len(files)}
        for filepath in files:
            try:
                doc_stats = self.process_document(filepath)
                stats["documents"].append(doc_stats)
                stats["total_chunks"] += doc_stats["chunk_count"]
            except Exception as e:
                print(f"  ERROR processing {filepath.name}: {e}")
                stats["documents"].append({
                    "filename": filepath.name, "error": str(e),
                })

        print(f"\n=== Done: {len(files)} documents, {stats['total_chunks']} chunks ===")
        return stats

    def process_document(self, filepath: Path) -> dict:
        from sqlalchemy import text
        t0 = time.monotonic()
        filename = filepath.name
        file_type = filepath.suffix.lower().lstrip(".")
        file_size = filepath.stat().st_size

        print(f"  Processing: {filename} ({file_size // 1024}K, {file_type})")

        # 1. Extract text
        markdown_text = self.processor.extract(filepath)
        print(f"    Extracted: {len(markdown_text)} chars")

        # 2. Chunk
        metadata = {"filename": filename, "file_type": file_type}
        chunks = self.chunker.chunk(markdown_text, metadata)
        print(f"    Chunks: {len(chunks)}")

        if not chunks:
            print(f"    WARNING: No chunks generated (file may be image-only)")
            return {"filename": filename, "chunk_count": 0, "warning": "no text extracted"}

        # 3. Embed
        chunk_texts = [c["content"] for c in chunks]
        embeddings = self.embedder.embed(chunk_texts)
        print(f"    Embedded: {len(embeddings)} vectors ({self.embedder.dimensions}d)")

        # 4. Store
        with self.pool.get_session() as session:
            # Upsert document
            result = session.execute(
                text("""
                    INSERT INTO polly_rag.documents (filename, file_type, title, file_size_bytes, markdown_text, chunk_count)
                    VALUES (:filename, :file_type, :title, :file_size, :markdown, :chunk_count)
                    ON CONFLICT (filename) DO UPDATE SET
                        markdown_text = EXCLUDED.markdown_text,
                        chunk_count = EXCLUDED.chunk_count,
                        processed_at = NOW()
                    RETURNING id
                """),
                {
                    "filename": filename,
                    "file_type": file_type,
                    "title": filepath.stem,
                    "file_size": file_size,
                    "markdown": markdown_text,
                    "chunk_count": len(chunks),
                },
            )
            doc_id = result.fetchone()[0]

            # Delete existing chunks for re-processing
            session.execute(
                text("DELETE FROM polly_rag.chunks WHERE document_id = :doc_id"),
                {"doc_id": doc_id},
            )

            # Insert chunks with embeddings
            for chunk, embedding in zip(chunks, embeddings):
                emb_str = "[" + ",".join(str(v) for v in embedding) + "]"
                session.execute(
                    text("""
                        INSERT INTO polly_rag.chunks (document_id, chunk_index, content, token_count, metadata, embedding)
                        VALUES (:doc_id, :idx, :content, :tokens, :meta, cast(:emb AS vector))
                    """),
                    {
                        "doc_id": doc_id,
                        "idx": chunk["chunk_index"],
                        "content": chunk["content"],
                        "tokens": chunk["token_count"],
                        "meta": json.dumps(chunk["metadata"]),
                        "emb": emb_str,
                    },
                )

        elapsed = time.monotonic() - t0
        print(f"    Stored: doc_id={doc_id}, {len(chunks)} chunks ({elapsed:.1f}s)")

        return {
            "filename": filename,
            "file_type": file_type,
            "doc_id": doc_id,
            "text_length": len(markdown_text),
            "chunk_count": len(chunks),
            "elapsed_seconds": round(elapsed, 1),
        }


def main():
    import argparse
    parser = argparse.ArgumentParser(description="Process documents into RAG vector store")
    parser.add_argument("--doc-dir", default=str(ROOT / "doc-data"),
                        help="Directory containing documents")
    args = parser.parse_args()

    doc_dir = Path(args.doc_dir)
    if not doc_dir.exists():
        print(f"ERROR: Document directory not found: {doc_dir}")
        sys.exit(1)

    indexer = RAGIndexer(doc_dir=doc_dir)
    stats = indexer.process_all()

    # Write stats
    stats_path = ROOT / "test-results" / "rag_indexing.json"
    stats_path.parent.mkdir(exist_ok=True)
    stats["timestamp"] = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    with open(stats_path, "w") as f:
        json.dump(stats, f, indent=2)
    print(f"\nStats written to {stats_path}")


if __name__ == "__main__":
    main()
